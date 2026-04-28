"""Build the thesis defense presentation as a .pptx file.

Big-picture, non-technical, ~10 min. Figures pulled from results/.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from PIL import Image

ROOT = Path("/home/tagesk/projects/GIN")
OUT = ROOT / "presentation.pptx"

NAVY = RGBColor(0x0F, 0x2E, 0x4F)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GRAY = RGBColor(0x55, 0x5C, 0x66)
LIGHT = RGBColor(0xF4, 0xF5, 0xF7)
DARK = RGBColor(0x1A, 0x1A, 0x1A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    bg.shadow.inherit = False
    return s


def add_title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.9))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.shadow.inherit = False

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), SW - Inches(1.0), Inches(0.7))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    r.font.name = "Calibri"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Calibri"
        r2.font.size = Pt(14)
        r2.font.color.rgb = RGBColor(0xCF, 0xD8, 0xE3)


def add_text(slide, x, y, w, h, lines, sizes=None, bolds=None, colors=None,
             align=PP_ALIGN.LEFT, line_spacing=1.15):
    """lines: list[str]. sizes/bolds/colors: parallel lists or scalars."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    n = len(lines)
    if sizes is None:
        sizes = [18] * n
    elif isinstance(sizes, int):
        sizes = [sizes] * n
    if bolds is None:
        bolds = [False] * n
    elif isinstance(bolds, bool):
        bolds = [bolds] * n
    if colors is None:
        colors = [DARK] * n
    elif not isinstance(colors, list):
        colors = [colors] * n

    for i, txt in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = txt
        r.font.name = "Calibri"
        r.font.size = Pt(sizes[i])
        r.font.bold = bolds[i]
        r.font.color.rgb = colors[i]
    return tb


def add_image(slide, path, x, y, max_w, max_h):
    """Place image preserving aspect, fitting inside (max_w, max_h)."""
    p = ROOT / path
    im = Image.open(p)
    iw, ih = im.size
    box_w = max_w
    box_h = max_h
    scale = min(box_w / Emu(iw * 9525), box_h / Emu(ih * 9525))
    # Just compute by ratio in EMU: pic native size in EMU = px * 9525 (assuming 96 dpi default).
    # Easier: compute aspect.
    ratio = iw / ih
    # Try filling width first
    w = max_w
    h = int(w / ratio)
    if h > max_h:
        h = max_h
        w = int(h * ratio)
    cx = x + (max_w - w) // 2
    cy = y + (max_h - h) // 2
    slide.shapes.add_picture(str(p), cx, cy, width=w, height=h)


def add_footer(slide, text):
    add_text(
        slide, Inches(0.5), Inches(7.15), SW - Inches(1.0), Inches(0.3),
        [text], sizes=10, colors=GRAY, align=PP_ALIGN.RIGHT, line_spacing=1.0,
    )


def add_math_box(slide, x, y, w, h, formula, accent=False, size=20):
    """Single-line math formula in a soft-tinted rounded box."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.line.color.rgb = ACCENT if accent else NAVY
    box.line.width = Pt(1)
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = formula
    r.font.name = "Cambria Math"
    r.font.size = Pt(size)
    r.font.italic = True
    r.font.color.rgb = ACCENT if accent else NAVY


# =====================================================================
# Slide 1: Title
# =====================================================================
s = add_slide()
# Big colored band
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.4), SW, Inches(2.7))
band.line.fill.background()
band.fill.solid()
band.fill.fore_color.rgb = NAVY
band.shadow.inherit = False

add_text(
    s, Inches(0.7), Inches(2.55), SW - Inches(1.4), Inches(1.4),
    ["Opening the Black Box of Graph Neural Networks"],
    sizes=40, bolds=True, colors=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.LEFT,
)
add_text(
    s, Inches(0.7), Inches(3.85), SW - Inches(1.4), Inches(1.0),
    ["Comparing 1-GNN and 1-2-GNN through generated explanations"],
    sizes=20, colors=RGBColor(0xCF, 0xD8, 0xE3), align=PP_ALIGN.LEFT,
)
add_text(
    s, Inches(0.7), Inches(5.6), SW - Inches(1.4), Inches(1.5),
    ["Master's thesis defense", "MUTAG  ·  PROTEINS"],
    sizes=[18, 14], colors=[NAVY, GRAY], align=PP_ALIGN.LEFT,
)

# =====================================================================
# Slide 2: Black-box motivation
# =====================================================================
s = add_slide()
add_title_bar(s, "Why this thesis?")

add_text(
    s, Inches(0.7), Inches(1.4), Inches(12), Inches(1.0),
    ["Neural networks make great decisions — but they don't tell us why."],
    sizes=24, bolds=True, colors=NAVY,
)
add_text(
    s, Inches(0.7), Inches(2.5), Inches(12), Inches(4.5),
    [
        "Two networks can reach the same accuracy on the same dataset…",
        "",
        "…and yet have completely different ideas of what each class actually looks like.",
        "",
        "If we only look at accuracy, we miss this. We never see what the model has actually learned.",
        "",
        "Goal: open the box. Force each model to show us its own picture of each class, then compare.",
    ],
    sizes=[20, 8, 20, 8, 20, 8, 22],
    bolds=[False, False, False, False, False, False, True],
    colors=[DARK, DARK, DARK, DARK, GRAY, DARK, ACCENT],
    line_spacing=1.2,
)

# =====================================================================
# Slide 3: What is a graph + a GNN
# =====================================================================
s = add_slide()
add_title_bar(s, "Graphs and Graph Neural Networks")

add_text(
    s, Inches(0.7), Inches(1.3), Inches(12), Inches(0.8),
    ["A graph = nodes connected by edges. Many real things are graphs."],
    sizes=22, bolds=True, colors=NAVY,
)

add_text(
    s, Inches(0.7), Inches(2.2), Inches(6), Inches(4.5),
    [
        "Examples in this thesis:",
        "",
        "• Molecules — atoms (nodes), bonds (edges)",
        "• Proteins — secondary-structure pieces, connected by sequence and 3-D distance",
        "",
        "Task: classify the whole graph.",
        "  – Is this molecule mutagenic?",
        "  – Is this protein an enzyme?",
    ],
    sizes=[20, 8, 18, 18, 8, 18, 16, 16],
    bolds=[True, False, False, False, False, True, False, False],
    line_spacing=1.25,
)

add_text(
    s, Inches(7.0), Inches(2.2), Inches(5.7), Inches(4.5),
    [
        "How a GNN works (intuition):",
        "",
        "• Every node holds a small vector",
        "• Each round, a node looks at its neighbours and updates its own vector",
        "• After a few rounds, all node vectors are pooled into a single graph vector",
        "• That graph vector is fed to a classifier",
    ],
    sizes=[20, 8, 17, 17, 17, 17],
    bolds=[True, False, False, False, False, False],
    line_spacing=1.25,
)

# =====================================================================
# Slide 4: 1-GNN vs 1-2-GNN
# =====================================================================
s = add_slide()
add_title_bar(s, "1-GNN vs. 1-2-GNN", "What does each architecture pay attention to?")

# Left column: 1-GNN
left_x = Inches(0.6)
right_x = Inches(7.0)
col_w = Inches(5.7)
add_text(
    s, left_x, Inches(1.3), col_w, Inches(0.7),
    ["1-GNN  —  the standard model"],
    sizes=22, bolds=True, colors=NAVY,
)
add_text(
    s, left_x, Inches(2.05), col_w, Inches(3.5),
    [
        "Each node only looks at its direct neighbours.",
        "",
        "• Cheap and fast",
        "• Powerful enough for most tasks",
        "• Limited: there are graph structures it provably cannot tell apart (e.g. two triangles vs. one hexagon — every node looks the same locally)",
    ],
    sizes=[18, 8, 17, 17, 17],
    bolds=[True, False, False, False, False],
    line_spacing=1.25,
)
# Math line — 1-GNN message passing
add_text(
    s, left_x, Inches(5.55), col_w, Inches(0.35),
    ["Message passing on nodes:"],
    sizes=13, bolds=True, colors=GRAY,
)
add_math_box(
    s, left_x, Inches(5.95), col_w, Inches(0.85),
    "h(v) ← σ( h(v) · W₁  +  Σᵤ∈N(v)  h(u) · W₂ )",
    accent=False, size=18,
)

# Right column: 1-2-GNN
add_text(
    s, right_x, Inches(1.3), col_w, Inches(0.7),
    ["1-2-GNN  —  also looks at pairs"],
    sizes=22, bolds=True, colors=ACCENT,
)
add_text(
    s, right_x, Inches(2.05), col_w, Inches(3.5),
    [
        "Does the 1-GNN step, then adds a second step over PAIRS of nodes.",
        "",
        "• A pair carries info about both nodes AND whether they are connected",
        "• Provably more expressive — can tell apart triangles vs. hexagons",
        "• Cost: pairs grow as N², so it is much slower and memory-hungry",
    ],
    sizes=[18, 8, 17, 17, 17],
    bolds=[True, False, False, False, False],
    line_spacing=1.25,
)
# Math line — 1-2-GNN message passing on pairs
add_text(
    s, right_x, Inches(5.55), col_w, Inches(0.35),
    ["Message passing on pairs:"],
    sizes=13, bolds=True, colors=GRAY,
)
add_math_box(
    s, right_x, Inches(5.95), col_w, Inches(0.85),
    "h(u,v) ← σ( h(u,v) · W₃  +  Σ pair-neighbours  h(·,·) · W₄ )",
    accent=True, size=17,
)

# =====================================================================
# Slide 5: The question
# =====================================================================
s = add_slide()
add_title_bar(s, "The question this thesis asks")

add_text(
    s, Inches(0.7), Inches(1.5), Inches(12), Inches(2),
    [
        "The 1-2-GNN can see more structure than the 1-GNN.",
        "On our datasets, both reach (almost) the same test accuracy.",
    ],
    sizes=22, line_spacing=1.3, colors=DARK,
)

add_text(
    s, Inches(0.7), Inches(3.6), Inches(12), Inches(1.5),
    ["So: are they actually learning the same thing?"],
    sizes=30, bolds=True, colors=ACCENT,
)

add_text(
    s, Inches(0.7), Inches(5.2), Inches(12), Inches(2),
    [
        "Same answer on a multiple-choice exam doesn't mean the two students reasoned the same way.",
        "Same accuracy on the test set doesn't mean the two models built the same idea of each class.",
    ],
    sizes=18, colors=GRAY, line_spacing=1.3,
)

# =====================================================================
# Slide 6: GIN-Graph idea
# =====================================================================
s = add_slide()
add_title_bar(s, "GIN-Graph: ask the model to draw its prototype")

add_text(
    s, Inches(0.7), Inches(1.15), Inches(12), Inches(0.6),
    ["The trick: reverse-engineer the model."],
    sizes=22, bolds=True, colors=NAVY,
)

# Noise → Generator → Graph flow (three boxes with arrows)
flow_y = Inches(1.95)
flow_h = Inches(1.05)
boxes = [
    ("Random noise", "z ~ 𝒩(0, I)", NAVY),
    ("Generator MLP", "trainable", ACCENT),
    ("Graph", "G̃ = (Ã, X̃)", NAVY),
]
box_w = Inches(2.7)
gap = Inches(0.6)
total = box_w * 3 + gap * 2
start_x = (SW - total) // 2
for i, (head, sub, col) in enumerate(boxes):
    bx = start_x + (box_w + gap) * i
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, flow_y, box_w, flow_h)
    box.line.color.rgb = col
    box.line.width = Pt(2)
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = head
    r.font.name = "Calibri"
    r.font.bold = True
    r.font.size = Pt(18)
    r.font.color.rgb = col
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = sub
    r2.font.name = "Cambria Math"
    r2.font.italic = True
    r2.font.size = Pt(14)
    r2.font.color.rgb = GRAY
    if i < 2:
        arrow_x = bx + box_w + Inches(0.05)
        arrow = s.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, arrow_x, flow_y + Inches(0.35),
            Inches(0.5), Inches(0.35),
        )
        arrow.line.fill.background()
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = NAVY
        arrow.shadow.inherit = False

# Generator math line
add_math_box(
    s, Inches(2.2), Inches(3.2), Inches(8.9), Inches(0.7),
    "G̃ = Generator(z) ,    z ~ 𝒩(0, I)",
    accent=False, size=18,
)

# Two reward bullets
add_text(
    s, Inches(0.7), Inches(4.1), Inches(12), Inches(1.6),
    [
        "The generator is rewarded when its graph G̃:",
        "    1.  is confidently classified as the target class by the frozen classifier",
        "    2.  looks realistic (a small adversary discriminator polices this)",
    ],
    sizes=[18, 17, 17],
    bolds=[True, False, False],
    line_spacing=1.3,
)

# Loss math line
add_text(
    s, Inches(0.7), Inches(5.7), Inches(12), Inches(0.35),
    ["Total generator loss:"],
    sizes=13, bolds=True, colors=GRAY, align=PP_ALIGN.CENTER,
)
add_math_box(
    s, Inches(2.2), Inches(6.05), Inches(8.9), Inches(0.75),
    "ℒ = (1−λ) · ℒ_GAN  +  λ · ℒ_GNN  +  ℒ_reg",
    accent=True, size=20,
)
add_text(
    s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.35),
    ["realism             class-specificity        size & node-type"],
    sizes=11, colors=GRAY, align=PP_ALIGN.CENTER,
)

# =====================================================================
# Slide 7: Pipeline
# =====================================================================
s = add_slide()
add_title_bar(s, "The pipeline")

steps = [
    ("1. Train classifier",
     "Train a 1-GNN and a 1-2-GNN on the dataset.\nEach one becomes a frozen “target”."),
    ("2. Train generator",
     "For every (model × class) pair, train a GIN-Graph generator.\n4 datasets × 2 classes × 2 models = 8 generators."),
    ("3. Sample graphs",
     "Generate 1000 graphs from each generator.\nThis is the model's prototype population."),
    ("4. Compare",
     "Compare against real data, and let each model classify the OTHER model's prototypes."),
]

start_y = Inches(1.4)
box_h = Inches(1.25)
gap = Inches(0.15)
for i, (head, body) in enumerate(steps):
    y = start_y + (box_h + gap) * i
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), y, Inches(2.6), box_h)
    box.line.fill.background()
    box.fill.solid()
    box.fill.fore_color.rgb = NAVY
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = head
    r.font.name = "Calibri"
    r.font.bold = True
    r.font.size = Pt(20)
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    add_text(
        s, Inches(3.4), y + Inches(0.1), Inches(9.5), box_h,
        body.split("\n"),
        sizes=17, line_spacing=1.2,
    )

# =====================================================================
# Slide 8: Datasets
# =====================================================================
s = add_slide()
add_title_bar(s, "Two datasets")

add_text(
    s, Inches(0.7), Inches(1.3), Inches(6), Inches(0.7),
    ["MUTAG  —  small molecules"],
    sizes=22, bolds=True, colors=NAVY,
)
add_text(
    s, Inches(0.7), Inches(2.0), Inches(6), Inches(3),
    [
        "• 188 molecules",
        "• 7 atom types (C, N, O, F, I, Cl, Br)",
        "• Classes: Mutagen / Non-Mutagen",
        "• Up to 28 atoms per graph",
    ],
    sizes=18, line_spacing=1.3,
)

add_text(
    s, Inches(7.0), Inches(1.3), Inches(6), Inches(0.7),
    ["PROTEINS  —  protein structures"],
    sizes=22, bolds=True, colors=NAVY,
)
add_text(
    s, Inches(7.0), Inches(2.0), Inches(6), Inches(3),
    [
        "• 1113 proteins",
        "• 3 secondary-structure types (Helix, Sheet, Coil/Turn)",
        "• Classes: Enzyme / Non-Enzyme",
        "• Up to 620 nodes; we cap generation at 50",
    ],
    sizes=18, line_spacing=1.3,
)

add_text(
    s, Inches(0.7), Inches(5.0), Inches(12), Inches(2),
    [
        "Test accuracy:",
        "    MUTAG       1-GNN  89.5 %     1-2-GNN  89.5 %",
        "    PROTEINS  1-GNN  79.8 %     1-2-GNN  78.9 %",
    ],
    sizes=[20, 18, 18],
    bolds=[True, False, False],
    colors=[ACCENT, DARK, DARK],
    line_spacing=1.4,
)

# =====================================================================
# Slide 9: MUTAG real vs generated samples
# =====================================================================
s = add_slide()
add_title_bar(s, "MUTAG  —  real vs. generated", "What does each model think a molecule looks like?")

# Three rows of wide images. Each is ~13" wide, 0.9" deep available
# Available area: y=1.0..7.1 → 6.1 inches. 3 rows + small label area.
labels_and_imgs = [
    ("Real molecules (from the dataset)",
     "results/mutag/comparison/figures/real_samples.png"),
    ("Generated by the 1-GNN",
     "results/mutag/comparison/figures/random_samples_1gnn.png"),
    ("Generated by the 1-2-GNN",
     "results/mutag/comparison/figures/random_samples_12gnn.png"),
]
row_y = [Inches(1.05), Inches(3.12), Inches(5.18)]
for (lbl, path), y in zip(labels_and_imgs, row_y):
    add_text(s, Inches(0.5), y, Inches(12), Inches(0.3),
             [lbl], sizes=14, bolds=True, colors=NAVY)
    add_image(s, path, Inches(0.5), y + Inches(0.3), Inches(12.3), Inches(1.7))

# =====================================================================
# Slide 10: PROTEINS real vs generated samples
# =====================================================================
s = add_slide()
add_title_bar(s, "PROTEINS  —  real vs. generated", "Same comparison, very different structure")

labels_and_imgs = [
    ("Real proteins (from the dataset)",
     "results/proteins/comparison/figures/real_samples.png"),
    ("Generated by the 1-GNN",
     "results/proteins/comparison/figures/random_samples_1gnn.png"),
    ("Generated by the 1-2-GNN",
     "results/proteins/comparison/figures/random_samples_12gnn.png"),
]
row_y = [Inches(1.05), Inches(3.12), Inches(5.18)]
for (lbl, path), y in zip(labels_and_imgs, row_y):
    add_text(s, Inches(0.5), y, Inches(12), Inches(0.3),
             [lbl], sizes=14, bolds=True, colors=NAVY)
    add_image(s, path, Inches(0.5), y + Inches(0.3), Inches(12.3), Inches(1.7))

# =====================================================================
# Slide 11: Cross-classification result
# =====================================================================
s = add_slide()
add_title_bar(s, "Headline result: cross-classification",
              "Does each model accept the OTHER model's prototypes?")

add_image(s, "results/mutag/comparison/figures/cross_classification.png",
          Inches(0.3), Inches(1.1), Inches(6.5), Inches(4.0))
add_image(s, "results/proteins/comparison/figures/cross_classification.png",
          Inches(6.7), Inches(1.1), Inches(6.5), Inches(4.0))

add_text(s, Inches(0.3), Inches(5.05), Inches(6.5), Inches(0.4),
         ["MUTAG"], sizes=14, bolds=True, colors=NAVY, align=PP_ALIGN.CENTER)
add_text(s, Inches(6.7), Inches(5.05), Inches(6.5), Inches(0.4),
         ["PROTEINS"], sizes=14, bolds=True, colors=NAVY, align=PP_ALIGN.CENTER)

add_text(
    s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(2),
    [
        "Each model accepts almost all of its OWN generated graphs.",
        "But in one specific class per dataset, the other model rejects them almost entirely:",
        "       MUTAG, Mutagen class:    1-GNN accepts only 12 % of the 1-2-GNN's prototypes.",
        "       PROTEINS, Enzyme class:  1-2-GNN accepts only 14 % of the 1-GNN's prototypes.",
    ],
    sizes=[18, 18, 17, 17],
    bolds=[False, True, False, False],
    colors=[DARK, ACCENT, DARK, DARK],
    line_spacing=1.25,
)

# =====================================================================
# Slide 12: Structural divergence (cycles + node-type shifts)
# =====================================================================
s = add_slide()
add_title_bar(s, "What drives the disagreement?",
              "MUTAG: triangles  ·  PROTEINS: too many sheets")

# Two side-by-side images: MUTAG cycles, PROTEINS node-type shifts
add_image(s, "results/mutag/comparison/figures/fingerprint_cycles.png",
          Inches(0.3), Inches(1.1), Inches(6.6), Inches(3.4))
add_image(s, "results/proteins/comparison/figures/fingerprint_node_type_shifts.png",
          Inches(7.0), Inches(1.1), Inches(6.0), Inches(3.4))

add_text(
    s, Inches(0.3), Inches(4.6), Inches(6.6), Inches(2.5),
    [
        "MUTAG  —  the theory-aligned signal",
        "The 1-2-GNN produces clearly fewer triangles than the 1-GNN.",
        "This is exactly the kind of structure that pair-level message passing can detect — and 1-WL cannot.",
    ],
    sizes=[16, 15, 14],
    bolds=[True, False, False],
    colors=[NAVY, DARK, GRAY],
    line_spacing=1.2,
)
add_text(
    s, Inches(7.0), Inches(4.6), Inches(6.0), Inches(2.5),
    [
        "PROTEINS  —  Sheet-heavy prototype",
        "The 1-2-GNN's Non-Enzyme prototype is ~89 % Sheet (real data: ~56 %).",
        "It exaggerates the structural pattern it has learned to associate with the class.",
    ],
    sizes=[16, 15, 14],
    bolds=[True, False, False],
    colors=[NAVY, DARK, GRAY],
    line_spacing=1.2,
)

# =====================================================================
# Slide 13: Conclusion
# =====================================================================
s = add_slide()
add_title_bar(s, "Conclusion")

add_text(
    s, Inches(0.7), Inches(1.3), Inches(12), Inches(1.5),
    ["Same accuracy does not mean the same understanding."],
    sizes=28, bolds=True, colors=ACCENT,
)

add_text(
    s, Inches(0.7), Inches(2.7), Inches(12), Inches(4.0),
    [
        "• On both MUTAG and PROTEINS the two architectures hit (almost) the same test accuracy.",
        "• Yet the prototypes they produce are systematically different — in atoms, edges, and cycles.",
        "• In one specific class per dataset, each model strongly rejects the other's prototypes.",
        "• The MUTAG triangle result is the cleanest theory-aligned signal: pair-level message passing changes which structures the model treats as evidence.",
        "",
        "Take-away:  higher-order message passing changes the prototype — not whether the model is more correct.",
        "Accuracy hides this. Generated explanations make it visible.",
    ],
    sizes=[18, 18, 18, 18, 8, 19, 18],
    bolds=[False, False, False, False, False, True, False],
    colors=[DARK, DARK, DARK, DARK, DARK, NAVY, GRAY],
    line_spacing=1.3,
)

add_text(
    s, Inches(0.7), Inches(6.8), Inches(12), Inches(0.5),
    ["Thank you  —  questions?"],
    sizes=20, bolds=True, colors=NAVY, align=PP_ALIGN.CENTER,
)

# =====================================================================
prs.save(OUT)
print(f"Saved {OUT}")
